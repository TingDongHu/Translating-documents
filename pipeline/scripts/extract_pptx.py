#!/usr/bin/env python3
"""extract_pptx.py — PowerPoint text extractor for translation pipelines.

Iterates all slides, all shapes (recursing into group shapes), and table cells.
Each paragraph gets a [Pnnn] tag with location metadata.
Group-nested shapes use a shape_path list instead of a flat shape index.
"""

import argparse
import json
import os


def parse_args():
    parser = argparse.ArgumentParser(
        description="Extract translatable text from a PPTX file."
    )
    parser.add_argument("--source", required=True)
    parser.add_argument("--extraction-output", required=True)
    parser.add_argument("--tagged-output", required=True)
    parser.add_argument("--report-output", required=True)
    return parser.parse_args()


def extract_from_shapes(shapes_collection, slide_idx, parent_path):
    """Recursively walk shapes (including group shapes) and extract text.

    Yields dicts with location including shape_path (list of indices
    for navigating nested group shapes).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape_idx, shape in enumerate(shapes_collection):
        current_path = parent_path + [shape_idx]

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # PYTHON-PPTX: group shapes have their own .shapes attribute
            # They are iterable even if not officially documented as such
            try:
                child_shapes = shape.shapes
                yield from extract_from_shapes(child_shapes, slide_idx, current_path)
            except Exception:
                pass
        else:
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if text:
                        yield {
                            "tag": None,  # assigned later
                            "type": "body",
                            "text": text,
                            "location": {
                                "slide": slide_idx,
                                "shape_path": current_path,
                                "para": para_idx,
                            },
                        }
            if shape.has_table:
                table = shape.table
                for row_idx in range(len(table.rows)):
                    for col_idx in range(len(table.columns)):
                        cell = table.cell(row_idx, col_idx)
                        text = cell.text.strip()
                        if text:
                            yield {
                                "tag": None,
                                "type": "table",
                                "text": text,
                                "location": {
                                    "slide": slide_idx,
                                    "shape_path": current_path,
                                    "row": row_idx,
                                    "col": col_idx,
                                },
                            }


def extract_from_pptx(source_path):
    """Extract text from all slide shapes, recursing into groups.

    Returns (paragraphs, warnings).
    """
    from pptx import Presentation

    paragraphs = []
    warnings = []

    prs = Presentation(source_path)
    tag_counter = 0

    for slide_idx, slide in enumerate(prs.slides):
        for item in extract_from_shapes(slide.shapes, slide_idx, []):
            item["tag"] = f"P{tag_counter}"
            tag_counter += 1
            paragraphs.append(item)

    # Also include a flat 'shape' field for backward compatibility
    # (maps the first element of shape_path for non-nested shapes)
    for p in paragraphs:
        loc = p["location"]
        if len(loc["shape_path"]) == 1:
            loc["shape"] = loc["shape_path"][0]

    return paragraphs, warnings


def ensure_parent_dir(filepath):
    parent = os.path.dirname(filepath)
    if parent:
        os.makedirs(parent, exist_ok=True)


def main():
    args = parse_args()
    paragraphs, warnings = extract_from_pptx(args.source)

    ensure_parent_dir(args.extraction_output)
    extraction = {"source": os.path.abspath(args.source), "paragraphs": paragraphs}
    with open(args.extraction_output, "w", encoding="utf-8") as f:
        json.dump(extraction, f, ensure_ascii=False, indent=2)

    ensure_parent_dir(args.tagged_output)
    with open(args.tagged_output, "w", encoding="utf-8") as f:
        for p in paragraphs:
            f.write(f"[{p['tag']}] {p['text']}\n")

    ensure_parent_dir(args.report_output)
    body_count = sum(1 for p in paragraphs if p['type'] == 'body')
    table_count = sum(1 for p in paragraphs if p['type'] == 'table')
    report = {
        "paragraph_count": len(paragraphs),
        "body_count": body_count,
        "table_count": table_count,
        "warnings": warnings,
    }
    with open(args.report_output, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)

    print(f"Extraction complete: {report['paragraph_count']} entries "
          f"({report['body_count']} body, {report['table_count']} table).")


if __name__ == "__main__":
    main()
