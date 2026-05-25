#!/usr/bin/env python3
"""Render translated text back into a PPTX presentation.

Reads an extraction.json, translated_merged.txt, and the original .pptx,
then writes a translated copy with text frame paragraphs and table cells
replaced. Supports nested group shapes via shape_path in locations.
"""

import argparse
import json
import re
from pathlib import Path

TAG_RE = re.compile(r'^\[([^\]]+)\]\s?(.*)$')


def load_translations(path):
    data = {}
    with open(path, encoding='utf-8') as f:
        for line in f:
            m = TAG_RE.match(line.rstrip('\n'))
            if m:
                data[m.group(1)] = m.group(2)
    return data


def resolve_shape(slide_or_group, shape_path):
    """Navigate into nested group shapes to find the target shape.

    shape_path is a list like [s0, s1, s2] where s0 is the top-level
    shape index, s1 is inside a group, s2 further nested, etc.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    target = None
    current_container = slide_or_group

    for depth, idx in enumerate(shape_path):
        shapes_collection = current_container.shapes
        target = shapes_collection[idx]
        # If there are more levels, the current shape must be a group
        if depth < len(shape_path) - 1:
            if target.shape_type == MSO_SHAPE_TYPE.GROUP:
                current_container = target
            else:
                # Path says go deeper but shape is not a group — stop
                return None

    return target


def build_lookup(paragraphs):
    """Build lookup dicts for text frames and table cells.

    Supports both flat 'shape' (backward compat) and 'shape_path' fields.
    """
    text_lookup = {}
    table_lookup = {}
    for p in paragraphs:
        loc = p['location']
        key = (loc['slide'], tuple(loc.get('shape_path', [])), loc.get('para'), loc.get('row'), loc.get('col'))
        if p['type'] == 'table':
            table_lookup[(loc['slide'], tuple(loc.get('shape_path', [])), loc.get('row'), loc.get('col'))] = p['tag']
        else:
            text_lookup[(loc['slide'], tuple(loc.get('shape_path', [])), loc.get('para'))] = p['tag']
    return text_lookup, table_lookup


def main():
    parser = argparse.ArgumentParser(
        description='Render translated text into a PPTX presentation.'
    )
    parser.add_argument('--template', required=True)
    parser.add_argument('--extraction', required=True)
    parser.add_argument('--translation', required=True)
    parser.add_argument('--output', required=True)
    parser.add_argument('--render-log', required=True)
    args = parser.parse_args()

    from pptx import Presentation

    translations = load_translations(args.translation)
    extraction = json.loads(Path(args.extraction).read_text(encoding='utf-8'))
    paragraphs = extraction.get('paragraphs', [])
    text_lookup, table_lookup = build_lookup(paragraphs)

    prs = Presentation(args.template)

    # Render text frame paragraphs
    for (slide_idx, shape_path, para_idx), tag in text_lookup.items():
        if tag not in translations:
            continue
        slide = prs.slides[slide_idx]
        shape = resolve_shape(slide, list(shape_path))
        if shape is None:
            continue
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if para_idx < len(tf.paragraphs):
            para = tf.paragraphs[para_idx]
            if para.runs:
                para.runs[0].text = translations[tag]
                for run in para.runs[1:]:
                    run.text = ''
            else:
                para.text = translations[tag]

    # Render table cells
    for (slide_idx, shape_path, row, col), tag in table_lookup.items():
        if tag not in translations:
            continue
        slide = prs.slides[slide_idx]
        shape = resolve_shape(slide, list(shape_path))
        if shape is None:
            continue
        if not shape.has_table:
            continue
        cell = shape.table.cell(row, col)
        tf = cell.text_frame
        for p in tf.paragraphs:
            p.text = ''
        tf.paragraphs[0].text = translations[tag]

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    applied = sum(1 for p in paragraphs if p['tag'] in translations)
    missing = [p['tag'] for p in paragraphs if p['tag'] not in translations]
    report = {
        'logical_tags_applied': applied,
        'missing': missing,
        'output_path': str(output_path),
        'layout_warnings': [],
    }
    log_path = Path(args.render_log)
    log_path.parent.mkdir(parents=True, exist_ok=True)
    log_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding='utf-8')
    print(f"Applied {applied} logical tags; missing {len(missing)}")


if __name__ == '__main__':
    main()
